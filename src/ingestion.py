import pandas as pd
import json
import pdfplumber
from pptx import Presentation

class DataIngestion:
    def __init__(self):
        self.data = []

    def read_csv(self, file_path):
        df = pd.read_csv(file_path)
        self.data.append(df)

    def read_json(self, file_path):
        with open(file_path, 'r') as f:
            json_data = json.load(f)
        df = pd.json_normalize(json_data)
        self.data.append(df)

    def read_pptx(self, file_path):
        prs = Presentation(file_path)
        slides_data = [shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")]
        self.data.append(pd.DataFrame(slides_data, columns=['Text']))

    def read_pdf(self, file_path):
        pdf_data = []
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                table = page.extract_table()
                if table:
                    pdf_data.extend(table)
        self.data.append(pd.DataFrame(pdf_data))

    def merge_data(self):
        return pd.concat(self.data, ignore_index=True)
